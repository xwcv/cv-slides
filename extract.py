#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
extract.py — 将《深度学习与计算机视觉》PPTX 课件抽取为中间 JSON + 图片资源。

用法：
    /usr/bin/python3 extract.py            # 抽取全部三讲
    /usr/bin/python3 extract.py 1          # 只抽取第 1 讲（可重复运行，幂等）

输出：
    data/lectureN.json   每页：title / bullets(带层级) / tables / images(路径+原始位置尺寸) / notes / skipped
    assets/lectureN/     图片按 s{页码:03d}_img{k}.{ext} 命名；同一 blob（sha1）只存一份，JSON 引用同一路径。

设计说明：
- 只读源 PPTX（OneDrive），绝不修改。
- emf/wmf 等浏览器不支持的格式不导出，但记录在 slides[i].skipped 中。
- SmartArt / 表格 / 图表：尽量提取文字。表格逐行归入 tables 字段；SmartArt/图表的 txBody 文字归入 bullets。
- 标题判定：优先 title placeholder；否则取最上方、字号最大的文本框首行。
"""
import sys
import os
import json
import hashlib
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

ROOT = os.path.dirname(os.path.abspath(__file__))
SRC_DIR = "/Users/xinggang/Library/CloudStorage/OneDrive-个人/Teaching/深度学习与计算机视觉-DLCV"

LECTURES = {
    1: {
        "file": "1 - 深度学习与计算机视觉简介.pptx",
        "name": "深度学习与计算机视觉简介",
        "slides": 79,
    },
    2: {
        "file": "2 - 从AlexNet到ChatGPT.pptx",
        "name": "从 AlexNet 到 ChatGPT",
        "slides": 98,
    },
    3: {
        "file": "3 - 从浅层到深度的表征学习.pptx",
        "name": "从浅层到深度的表征学习",
        "slides": 112,
    },
}

# 浏览器可直接显示的图片格式；其余跳过但记录
OK_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".webp", ".tif", ".tiff"}


def clean_text(s):
    """压缩空白，去掉纯空内容。"""
    if s is None:
        return ""
    s = s.replace("\x0b", "\n").replace("\r", "\n")
    lines = [re.sub(r"\s+", " ", ln).strip() for ln in s.split("\n")]
    return "\n".join([ln for ln in lines if ln])


def iter_text_frame(tf, base_level=0):
    """把 text_frame 的段落转成 [(level, text), ...]，去掉空项。"""
    out = []
    for para in tf.paragraphs:
        text = clean_text("".join(run.text for run in para.runs) or para.text)
        if not text:
            continue
        level = base_level + (para.level or 0)
        out.append((level, text))
    return out


def shape_text(shape):
    """提取单个 shape 的全部文字段落（含 group 递归）。"""
    items = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            items.extend(shape_text(sub))
        return items
    if shape.has_text_frame:
        items.extend(iter_text_frame(shape.text_frame))
    return items


def table_text(graphic_frame):
    """表格逐行提取，每行内单元格以 ' | ' 连接。"""
    rows = []
    try:
        tbl = graphic_frame.table
    except Exception:
        return rows
    for row in tbl.rows:
        cells = [clean_text(c.text) for c in row.cells]
        rows.append(cells)
    return rows


def max_font_size(text_frame):
    """文本框中显式设置过的最大字号（pt），未设置返回 None。"""
    best = None
    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                pt = run.font.size.pt
                if best is None or pt > best:
                    best = pt
    return best


def extract_slide(slide, slide_idx, img_store, assets_rel, slide_h_emu=0):
    """抽取单页 -> dict。img_store: {sha1: relpath}，跨页去重。"""
    title = ""
    title_source = False
    bullets = []          # [(level, text)]
    tables = []           # [[row...], ...]
    images = []           # [{path,left,top,width,height,rot}]
    skipped = []          # [{kind, name, detail}]

    # 1) title placeholder
    try:
        if slide.shapes.title is not None:
            t = clean_text(slide.shapes.title.text)
            if t:
                title = t.split("\n")[0]
                title_source = True
    except Exception:
        pass

    text_shapes = []  # (shape, items, top, max_pt) 供标题推断

    def walk(shapes):
        for shape in shapes:
            st = shape.shape_type
            try:
                name = shape.name or ""
            except Exception:
                name = ""

            if st == MSO_SHAPE_TYPE.GROUP:
                try:
                    walk(shape.shapes)
                except Exception:
                    pass
                continue

            if st == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    ext = os.path.splitext(shape.image.filename or "")[1].lower()
                    if not ext:
                        ext = "." + (shape.image.ext or "png")
                except Exception as e:
                    skipped.append({"kind": "picture", "name": name,
                                    "detail": "读取失败: %s" % e})
                    continue
                sha1 = hashlib.sha1(blob).hexdigest()
                if ext not in OK_EXT:
                    skipped.append({"kind": "picture", "name": name,
                                    "detail": "浏览器不支持的格式 %s（sha1=%s）" % (ext, sha1[:10])})
                    continue
                if sha1 not in img_store:
                    fname = "s%03d_img%d%s" % (slide_idx, len(img_store) + 1, ext)
                    fpath = os.path.join(ROOT, assets_rel, fname)
                    with open(fpath, "wb") as f:
                        f.write(blob)
                    img_store[sha1] = "%s/%s" % (assets_rel, fname)
                images.append({
                    "path": img_store[sha1],
                    "left": int(shape.left or 0), "top": int(shape.top or 0),
                    "width": int(shape.width or 0), "height": int(shape.height or 0),
                    "rot": float(shape.rotation or 0),
                })
                continue

            if st == MSO_SHAPE_TYPE.TABLE or (getattr(shape, "has_table", False)):
                rows = table_text(shape)
                if rows:
                    tables.append(rows)
                else:
                    skipped.append({"kind": "table", "name": name, "detail": "空表格"})
                continue

            if getattr(shape, "has_chart", False):
                # 图表：尝试取标题/分类文字
                txt = []
                try:
                    ch = shape.chart
                    if ch.has_title:
                        txt.append(clean_text(ch.chart_title.text_frame.text))
                    for plot in ch.plots:
                        txt.extend([str(c) for c in plot.categories])
                except Exception:
                    pass
                txt = [t for t in txt if t]
                if txt:
                    bullets.extend((1, t) for t in txt)
                skipped.append({"kind": "chart", "name": name,
                                "detail": "图表仅提取文字，图形未转换"})
                continue

            if shape.has_text_frame:
                items = iter_text_frame(shape.text_frame)
                if items:
                    pt = max_font_size(shape.text_frame)
                    text_shapes.append((shape, items,
                                        int(shape.top or 0), pt or 0))
                continue

            # SmartArt（diagram / IGX_GRAPHIC）/ 其他图形：尽量提取文字
            try:
                txt = shape_text(shape)
                if txt:
                    bullets.extend(txt)
                elif st != MSO_SHAPE_TYPE.LINE:  # 装饰线条不计入 skipped
                    skipped.append({"kind": str(st), "name": name,
                                    "detail": "无可提取文字"})
            except Exception:
                skipped.append({"kind": str(st), "name": name,
                                "detail": "解析失败"})
            continue

    walk(slide.shapes)

    # 2) 标题推断：无 title placeholder 时，在页面上半部分找字号最大且不像日期/页码的文本框首行
    used_as_title = None
    DATE_RE = re.compile(r"^[\d\-\./年月日\s:A-Za-z]{4,30}$")
    if not title and text_shapes:
        slide_h_limit = int(slide_h_emu * 0.5) if slide_h_emu else 3_000_000
        cands = []
        for shape, items, top, pt in text_shapes:
            first = items[0][1]
            if len(first) > 60 or len(first) < 2:
                continue
            if DATE_RE.match(first) and any(c.isdigit() for c in first):
                continue  # 日期 / 页码 / 纯编号
            if top <= slide_h_limit:
                cands.append((shape, items, top, pt))
        if cands:
            max_pt = max(c[3] for c in cands)
            # 字号并列时取最靠上的
            best = sorted([c for c in cands if c[3] >= max_pt - 2],
                          key=lambda x: x[2])[0]
            title = best[1][0][1]
            used_as_title = best

    # 3) 其余文本按 top 排序归入 bullets；标题所在 shape 的首行若与 title 相同则跳过
    for shape, items, top, pt in sorted(text_shapes, key=lambda x: x[2]):
        if shape is not None and title and used_as_title is not None and shape is used_as_title[0]:
            items = items[1:]
        # 跳过与 title placeholder 完全相同的文本框（模板常见重复）
        items = [(lv, tx) for (lv, tx) in items if tx != title]
        bullets.extend(items)

    # 去重相邻完全相同的 bullet
    deduped = []
    for it in bullets:
        if not deduped or deduped[-1] != it:
            deduped.append(it)
    bullets = deduped

    # 4) 演讲者备注
    notes = ""
    try:
        if slide.has_notes_slide:
            notes = clean_text(slide.notes_slide.notes_text_frame.text)
    except Exception:
        pass

    return {
        "index": slide_idx,
        "title": title,
        "bullets": [{"level": lv, "text": tx} for (lv, tx) in bullets],
        "tables": tables,
        "images": images,
        "notes": notes,
        "skipped": skipped,
    }


def extract_lecture(num):
    info = LECTURES[num]
    src = os.path.join(SRC_DIR, info["file"])
    assets_rel = "assets/lecture%d" % num
    os.makedirs(os.path.join(ROOT, assets_rel), exist_ok=True)
    os.makedirs(os.path.join(ROOT, "data"), exist_ok=True)

    prs = Presentation(src)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    img_store = {}
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        slides.append(extract_slide(slide, i, img_store, assets_rel, slide_h))

    doc = {
        "lecture": num,
        "name": info["name"],
        "source": info["file"],
        "slide_width_emu": slide_w,
        "slide_height_emu": slide_h,
        "slide_count": len(slides),
        "image_count": len(img_store),
        "slides": slides,
    }
    out = os.path.join(ROOT, "data", "lecture%d.json" % num)
    with open(out, "w", encoding="utf-8") as f:
        json.dump(doc, f, ensure_ascii=False, indent=1)

    n_img_ref = sum(len(s["images"]) for s in slides)
    n_skip = sum(len(s["skipped"]) for s in slides)
    print("lecture%d: %d 页, 图片引用 %d 处（去重后 %d 个文件）, 跳过元素 %d 个 -> %s"
          % (num, len(slides), n_img_ref, len(img_store), n_skip, out))
    return doc


def main():
    nums = [int(a) for a in sys.argv[1:]] or sorted(LECTURES)
    for n in nums:
        extract_lecture(n)


if __name__ == "__main__":
    main()
